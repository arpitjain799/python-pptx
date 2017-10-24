# encoding: utf-8

"""Unit-test suite for pptx.shapes.freeform module"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.shapes.autoshape import Shape
from pptx.shapes.freeform import (
    _BaseDrawingOperation, _Close, FreeformBuilder, _LineSegment
)
from pptx.shapes.shapetree import SlideShapes

from ..unitutil.cxml import element
from ..unitutil.file import snippet_seq
from ..unitutil.mock import (
    call, initializer_mock, instance_mock, method_mock, property_mock
)


class DescribeFreeformBuilder(object):

    def it_provides_a_constructor(self, new_fixture):
        shapes_, start_x, start_y, x_scale, y_scale = new_fixture[:5]
        _init_, start_x_int, start_y_int = new_fixture[5:]

        builder = FreeformBuilder.new(
            shapes_, start_x, start_y, x_scale, y_scale
        )

        _init_.assert_called_once_with(
            builder, shapes_, start_x_int, start_y_int, x_scale, y_scale
        )
        assert isinstance(builder, FreeformBuilder)

    def it_can_add_straight_line_segments(self, add_segs_fixture):
        builder, vertices, close, add_calls, close_calls = add_segs_fixture

        return_value = builder.add_line_segments(vertices, close)

        assert builder._add_line_segment.call_args_list == add_calls
        assert builder._add_close.call_args_list == close_calls
        assert return_value is builder

    def it_can_build_the_specified_freeform_shape(self, convert_fixture):
        builder, origin_x, origin_y, sp = convert_fixture[:4]
        apply_operation_to_, calls, shape_ = convert_fixture[4:]

        shape = builder.convert_to_shape(origin_x, origin_y)

        builder._add_freeform_sp.assert_called_once_with(
            builder, origin_x, origin_y
        )
        builder._start_path.assert_called_once_with(builder, sp)
        assert apply_operation_to_.call_args_list == calls
        builder._shapes._shape_factory.assert_called_once_with(sp)
        assert shape is shape_

    def it_knows_the_shape_x_offset(self, shape_offset_x_fixture):
        builder, expected_value = shape_offset_x_fixture
        x_offset = builder.shape_offset_x
        assert x_offset == expected_value

    def it_knows_the_shape_y_offset(self, shape_offset_y_fixture):
        builder, expected_value = shape_offset_y_fixture
        y_offset = builder.shape_offset_y
        assert y_offset == expected_value

    def it_adds_a_freeform_sp_to_help(self, sp_fixture):
        builder, origin_x, origin_y, spTree, expected_xml = sp_fixture

        sp = builder._add_freeform_sp(origin_x, origin_y)

        assert spTree.xml == expected_xml
        assert sp is spTree.xpath('p:sp')[0]

    def it_adds_a_line_segment_to_help(self, add_seg_fixture):
        builder, x, y, _LineSegment_new_, line_segment_ = add_seg_fixture

        builder._add_line_segment(x, y)

        _LineSegment_new_.assert_called_once_with(builder, x, y)
        assert builder._drawing_operations == [line_segment_]

    def it_closes_a_contour_to_help(self, add_close_fixture):
        builder, _Close_new_, close_ = add_close_fixture

        builder._add_close()

        _Close_new_.assert_called_once_with()
        assert builder._drawing_operations == [close_]

    def it_knows_the_freeform_left_extent_to_help(self, left_fixture):
        builder, expected_value = left_fixture
        left = builder._left
        assert left == expected_value

    def it_knows_the_freeform_top_extent_to_help(self, top_fixture):
        builder, expected_value = top_fixture
        top = builder._top
        assert top == expected_value

    def it_knows_the_freeform_width_to_help(self, width_fixture):
        builder, expected_value = width_fixture
        width = builder._width
        assert width == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def add_close_fixture(self, _Close_new_, close_):
        _Close_new_.return_value = close_
        builder = FreeformBuilder(None, None, None, None, None)
        return builder, _Close_new_, close_

    @pytest.fixture
    def add_seg_fixture(self, _LineSegment_new_, line_segment_):
        x, y = 4, 2
        _LineSegment_new_.return_value = line_segment_

        builder = FreeformBuilder(None, None, None, None, None)
        return builder, x, y, _LineSegment_new_, line_segment_

    @pytest.fixture(params=[
        (True,  [call()]),
        (False, []),
    ])
    def add_segs_fixture(self, request, _add_line_segment_, _add_close_):
        close, close_calls = request.param
        vertices = ((1, 2), (3, 4), (5, 6))
        builder = FreeformBuilder(None, None, None, None, None)
        add_calls = [call(1, 2), call(3, 4), call(5, 6)]
        return builder, vertices, close, add_calls, close_calls

    @pytest.fixture
    def convert_fixture(self, shapes_, apply_operation_to_,
                        _add_freeform_sp_, _start_path_, shape_):
        origin_x, origin_y = 42, 24
        sp, path = element('p:sp'), element('a:path')
        drawing_ops = (
            _BaseDrawingOperation(None, None, None),
            _BaseDrawingOperation(None, None, None),
        )
        shapes_._shape_factory.return_value = shape_
        _add_freeform_sp_.return_value = sp
        _start_path_.return_value = path

        builder = FreeformBuilder(shapes_, None, None, None, None)
        builder._drawing_operations.extend(drawing_ops)
        calls = [call(drawing_ops[0], path), call(drawing_ops[1], path)]
        return (
            builder, origin_x, origin_y, sp, apply_operation_to_, calls,
            shape_
        )

    @pytest.fixture(params=[
        (0,      1.0,   0),
        (4,      10.0,  40),
        (914400, 914.3, 836035920),
    ])
    def left_fixture(self, request, shape_offset_x_prop_):
        offset_x, x_scale, expected_value = request.param
        shape_offset_x_prop_.return_value = offset_x

        builder = FreeformBuilder(None, None, None, x_scale, None)
        return builder, expected_value

    @pytest.fixture
    def new_fixture(self, shapes_, _init_):
        start_x, start_y, x_scale, y_scale = 99.56, 200.49, 4.2, 2.4
        start_x_int, start_y_int = 100, 200
        return (
            shapes_, start_x, start_y, x_scale, y_scale, _init_, start_x_int,
            start_y_int
        )

    @pytest.fixture(params=[
        (0,  (1, None, 2, 3),       0),
        (6,  (1, None, 2, 3),       1),
        (50, (150, -5, None, 100), -5),
    ])
    def shape_offset_x_fixture(self, request):
        start_x, xs, expected_value = request.param
        drawing_ops = []
        for x in xs:
            if x is None:
                drawing_ops.append(_Close())
            else:
                drawing_ops.append(_BaseDrawingOperation(None, x, None))

        builder = FreeformBuilder(None, start_x, None, None, None)
        builder._drawing_operations.extend(drawing_ops)
        return builder, expected_value

    @pytest.fixture(params=[
        (0,  (2, None, 6, 8),        0),
        (4,  (2, None, 6, 8),        2),
        (19, (213, -22, None, 100), -22),
    ])
    def shape_offset_y_fixture(self, request):
        start_y, ys, expected_value = request.param
        drawing_ops = []
        for y in ys:
            if y is None:
                drawing_ops.append(_Close())
            else:
                drawing_ops.append(_BaseDrawingOperation(None, None, y))

        builder = FreeformBuilder(None, None, start_y, None, None)
        builder._drawing_operations.extend(drawing_ops)
        return builder, expected_value

    @pytest.fixture
    def sp_fixture(self, _left_prop_, _top_prop_, _width_prop_,
                   _height_prop_):
        origin_x, origin_y = 42, 24
        spTree = element('p:spTree')
        shapes = SlideShapes(spTree, None)
        _left_prop_.return_value, _top_prop_.return_value = 12, 34
        _width_prop_.return_value, _height_prop_.return_value = 56, 78

        builder = FreeformBuilder(shapes, None, None, None, None)
        expected_xml = snippet_seq('freeform')[0]
        return builder, origin_x, origin_y, spTree, expected_xml

    @pytest.fixture(params=[
        (0,      11.0,  0),
        (100,    10.36, 1036),
        (914242, 943.1, 862221630),
    ])
    def top_fixture(self, request, shape_offset_y_prop_):
        offset_y, y_scale, expected_value = request.param
        shape_offset_y_prop_.return_value = offset_y

        builder = FreeformBuilder(None, None, None, None, y_scale)
        return builder, expected_value

    @pytest.fixture(params=[
        (0,      1.0,   0),
        (42,     10.0,  420),
        (914400, 914.4, 836127360),
    ])
    def width_fixture(self, request, _dx_prop_):
        dx, x_scale, expected_value = request.param
        _dx_prop_.return_value = dx

        builder = FreeformBuilder(None, None, None, x_scale, None)
        return builder, expected_value

    # fixture components -----------------------------------

    @pytest.fixture
    def _add_close_(self, request):
        return method_mock(request, FreeformBuilder, '_add_close')

    @pytest.fixture
    def _add_freeform_sp_(self, request):
        return method_mock(
            request, FreeformBuilder, '_add_freeform_sp', autospec=True
        )

    @pytest.fixture
    def _add_line_segment_(self, request):
        return method_mock(request, FreeformBuilder, '_add_line_segment')

    @pytest.fixture
    def apply_operation_to_(self, request):
        return method_mock(
            request, _BaseDrawingOperation, 'apply_operation_to',
            autospec=True
        )

    @pytest.fixture
    def close_(self, request):
        return instance_mock(request, _Close)

    @pytest.fixture
    def _Close_new_(self, request):
        return method_mock(request, _Close, 'new')

    @pytest.fixture
    def _dx_prop_(self, request):
        return property_mock(request, FreeformBuilder, '_dx')

    @pytest.fixture
    def _height_prop_(self, request):
        return property_mock(request, FreeformBuilder, '_height')

    @pytest.fixture
    def _init_(self, request):
        return initializer_mock(request, FreeformBuilder, autospec=True)

    @pytest.fixture
    def _left_prop_(self, request):
        return property_mock(request, FreeformBuilder, '_left')

    @pytest.fixture
    def line_segment_(self, request):
        return instance_mock(request, _LineSegment)

    @pytest.fixture
    def _LineSegment_new_(self, request):
        return method_mock(request, _LineSegment, 'new')

    @pytest.fixture
    def shape_(self, request):
        return instance_mock(request, Shape)

    @pytest.fixture
    def shape_offset_x_prop_(self, request):
        return property_mock(request, FreeformBuilder, 'shape_offset_x')

    @pytest.fixture
    def shape_offset_y_prop_(self, request):
        return property_mock(request, FreeformBuilder, 'shape_offset_y')

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, SlideShapes)

    @pytest.fixture
    def _start_path_(self, request):
        return method_mock(
            request, FreeformBuilder, '_start_path', autospec=True
        )

    @pytest.fixture
    def _top_prop_(self, request):
        return property_mock(request, FreeformBuilder, '_top')

    @pytest.fixture
    def _width_prop_(self, request):
        return property_mock(request, FreeformBuilder, '_width')


class Describe_BaseDrawingOperation(object):

    def it_knows_its_x_coordinate(self, x_fixture):
        drawing_operation, expected_value = x_fixture
        x = drawing_operation.x
        assert x == expected_value

    def it_knows_its_y_coordinate(self, y_fixture):
        drawing_operation, expected_value = y_fixture
        y = drawing_operation.y
        assert y == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def x_fixture(self):
        x = 42
        drawing_operation = _BaseDrawingOperation(None, x, None)
        expected_value = x
        return drawing_operation, expected_value

    @pytest.fixture
    def y_fixture(self):
        y = 24
        drawing_operation = _BaseDrawingOperation(None, None, y)
        expected_value = y
        return drawing_operation, expected_value


class Describe_Close(object):

    def it_provides_a_constructor(self, new_fixture):
        _init_ = new_fixture

        close = _Close.new()

        _init_.assert_called_once_with()
        assert isinstance(close, _Close)

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def new_fixture(self, _init_):
        return _init_

    # fixture components -----------------------------------

    @pytest.fixture
    def _init_(self, request):
        return initializer_mock(request, _Close, autospec=True)


class Describe_LineSegment(object):

    def it_provides_a_constructor(self, new_fixture):
        builder_, x, y, _init_, x_int, y_int = new_fixture

        line_segment = _LineSegment.new(builder_, x, y)

        _init_.assert_called_once_with(line_segment, builder_, x_int, y_int)
        assert isinstance(line_segment, _LineSegment)

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def new_fixture(self, builder_, _init_):
        x, y, x_int, y_int = 99.51, 200.49, 100, 200
        return builder_, x, y, _init_, x_int, y_int

    # fixture components -----------------------------------

    @pytest.fixture
    def builder_(self, request):
        return instance_mock(request, FreeformBuilder)

    @pytest.fixture
    def _init_(self, request):
        return initializer_mock(request, _LineSegment, autospec=True)